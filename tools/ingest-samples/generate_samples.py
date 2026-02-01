from pathlib import Path
import json
import csv
from datetime import date

from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import Workbook


BASE_DIR = Path(__file__).resolve().parent
BASE_DIR.mkdir(parents=True, exist_ok=True)

today = date.today().isoformat()
meta_block = (
    f"DOC_ID: SAMPLE-001\n"
    f"SECTOR: ENTERPRISE\n"
    f"TITLE: Sample Ingestion File\n"
    f"DATE: {today}\n"
    f"AUTHOR_ROLE: Operations Analyst\n"
    f"CLASSIFICATION: UNCLASSIFIED // FOR TRAINING\n"
    f"SUMMARY: Sample content for ingestion validation.\n\n"
)


def write_text(filename: str, body: str):
    (BASE_DIR / filename).write_text(meta_block + body, encoding="utf-8")


def write_json(filename: str, data: dict):
    payload = {
        "metadata": {
            "doc_id": "SAMPLE-JSON-001",
            "sector": "ENTERPRISE",
            "title": "Sample JSON Ingestion File",
            "date": today,
            "author_role": "Data Engineer",
            "classification": "UNCLASSIFIED // FOR TRAINING",
        },
        "summary": "Sample JSON content for ingestion validation.",
        "data": data,
    }
    (BASE_DIR / filename).write_text(json.dumps(payload, indent=2), encoding="utf-8")


def write_ndjson(filename: str, rows: list[dict]):
    with (BASE_DIR / filename).open("w", encoding="utf-8") as f:
        f.write(json.dumps({"metadata": {"doc_id": "SAMPLE-NDJSON-001", "sector": "ENTERPRISE"}}) + "\n")
        for row in rows:
            f.write(json.dumps(row) + "\n")


def write_csv(filename: str, rows: list[dict]):
    with (BASE_DIR / filename).open("w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=list(rows[0].keys()))
        writer.writeheader()
        writer.writerows(rows)


def write_docx(filename: str):
    doc = DocxDocument()
    doc.add_heading("Sample DOCX Ingestion File", level=1)
    doc.add_paragraph(meta_block.strip())
    doc.add_paragraph("This DOCX is used to validate ingestion of policy memos and procedures.")
    doc.add_paragraph("Key Findings:\n- Control effectiveness improved by 18%\n- Incident response time reduced by 22%\n")
    doc.save(BASE_DIR / filename)


def write_pptx(filename: str):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Sample PPTX Ingestion File"
    slide.placeholders[1].text = (
        f"SECTOR: ENTERPRISE\nDATE: {today}\n"
        "Briefing highlights:\n- Program status: Green\n- Budget burn: 41%\n- Risks: 2 open\n"
    )
    prs.save(BASE_DIR / filename)


def write_xlsx(filename: str):
    wb = Workbook()
    ws = wb.active
    ws.title = "Metrics"
    ws.append(["Metric", "Value", "Notes"])
    ws.append(["System Availability", "99.4%", "Quarterly average"])
    ws.append(["Incident Count", 12, "Last 90 days"])
    ws.append(["Budget Utilization", "71%", "FY to date"])
    wb.save(BASE_DIR / filename)


def main():
    write_text("sample.txt", "This is a plain text ingestion sample with operational notes.\n")
    write_text("sample.md", "# Sample Markdown\n\n- Item A\n- Item B\n")
    write_text("sample.log", "2026-01-30 14:00:00 INFO Service started\n2026-01-30 14:05:12 WARN Threshold exceeded\n")
    write_text("sample.html", "<html><body><h1>Sample HTML</h1><p>Internal portal report.</p></body></html>")

    write_csv(
        "sample.csv",
        [
            {"Metric": "Throughput", "Value": "1200 rpm", "Owner": "Ops"},
            {"Metric": "Error Rate", "Value": "0.7%", "Owner": "SRE"},
        ],
    )

    write_json(
        "sample.json",
        {
            "systems": [{"name": "Edge Cluster", "status": "green"}],
            "risks": [{"id": "R-102", "severity": "medium"}],
        },
    )

    write_ndjson(
        "sample.ndjson",
        [
            {"timestamp": today, "event": "job_started", "job_id": "JOB-001"},
            {"timestamp": today, "event": "job_completed", "job_id": "JOB-001"},
        ],
    )

    write_docx("sample.docx")
    write_pptx("sample.pptx")
    write_xlsx("sample.xlsx")


if __name__ == "__main__":
    main()
