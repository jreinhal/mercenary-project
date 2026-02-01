from __future__ import annotations

from pathlib import Path
import json
import random
from datetime import datetime, timezone

from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import Workbook


BASE = Path(r"D:\Corpus")
SECTORS = ["Government", "Medical", "Finance", "Academic", "Enterprise"]

# Additional formats to reach ~23% non-txt total per sector
FORMAT_COUNTS = {
    ".docx": 20,
    ".pptx": 15,
    ".xlsx": 15,
    ".html": 20,
    ".json": 10,
    ".ndjson": 10,
    ".log": 10,
}

random.seed(42)


def read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def write_docx(path: Path, content: str) -> None:
    doc = DocxDocument()
    doc.add_heading(path.stem.replace("_", " ").title(), level=1)
    for para in content.splitlines()[:80]:
        if para.strip():
            doc.add_paragraph(para.strip())
    doc.save(path)


def write_pptx(path: Path, content: str) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = path.stem.replace("_", " ").title()
    bullets = [line.strip() for line in content.splitlines() if line.strip()][:6]
    slide.placeholders[1].text = "\n".join(bullets) if bullets else "Briefing summary."
    prs.save(path)


def write_xlsx(path: Path, content: str) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Line"])
    for line in content.splitlines()[:200]:
        if line.strip():
            ws.append([line.strip()])
    wb.save(path)


def write_html(path: Path, content: str) -> None:
    html = "<html><body><pre>" + content + "</pre></body></html>"
    path.write_text(html, encoding="utf-8")


def write_json(path: Path, content: str) -> None:
    payload = {
        "doc_id": path.stem,
        "created_at": datetime.now(timezone.utc).isoformat().replace("+00:00", "Z"),
        "content": content.splitlines()[:120],
    }
    path.write_text(json.dumps(payload, indent=2), encoding="utf-8")


def write_ndjson(path: Path, content: str) -> None:
    lines = [line.strip() for line in content.splitlines() if line.strip()][:120]
    with path.open("w", encoding="utf-8") as f:
        f.write(json.dumps({"doc_id": path.stem}) + "\n")
        for line in lines:
            f.write(json.dumps({"line": line}) + "\n")


def write_log(path: Path, content: str) -> None:
    lines = [line.strip() for line in content.splitlines() if line.strip()][:120]
    now = datetime.now(timezone.utc)
    formatted = []
    for i, line in enumerate(lines):
        ts = now.replace(microsecond=0).isoformat() + "Z"
        formatted.append(f"{ts} INFO {line}")
    path.write_text("\n".join(formatted), encoding="utf-8")


def convert_file(src: Path, ext: str) -> None:
    content = read_text(src)
    dest = src.with_suffix(ext)
    if dest.exists():
        return
    if ext == ".docx":
        write_docx(dest, content)
    elif ext == ".pptx":
        write_pptx(dest, content)
    elif ext == ".xlsx":
        write_xlsx(dest, content)
    elif ext == ".html":
        write_html(dest, content)
    elif ext == ".json":
        write_json(dest, content)
    elif ext == ".ndjson":
        write_ndjson(dest, content)
    elif ext == ".log":
        write_log(dest, content)


def main() -> None:
    for sector in SECTORS:
        folder = BASE / sector
        if not folder.exists():
            continue
        txt_files = list(folder.glob("*.txt"))
        existing_variants = {p.stem for p in folder.glob("*.md")} | {p.stem for p in folder.glob("*.csv")}
        candidates = [p for p in txt_files if p.stem not in existing_variants]
        random.shuffle(candidates)

        cursor = 0
        for ext, count in FORMAT_COUNTS.items():
            batch = candidates[cursor: cursor + count]
            cursor += count
            for src in batch:
                convert_file(src, ext)


if __name__ == "__main__":
    main()
