#!/usr/bin/env python3
"""
Enterprise Corpus Generator for Sentinel RAG Testing.

Generates realistic, template-based documents across GOVERNMENT, MEDICAL, and
ENTERPRISE sectors.  Zero LLM tokens — all content is deterministic templates
with randomised slot values (names, dates, dollar amounts, etc.).

Additionally generates deterministic "golden" enterprise documents with fixed
facts that answer specific test queries from the Enterprise Test Plan.

Usage:
    python generate_corpus.py                              # 3000 docs to D:\\Corpus
    python generate_corpus.py --output D:\\Corpus --count 3000
    python generate_corpus.py --count 600                  # 200 per sector
    python generate_corpus.py --golden-only                # Only generate golden docs

Supported output formats (weighted distribution):
    .txt  40%  |  .docx 12%  |  .xlsx/.csv 10%  |  .html 8%
    .md    5%  |  .pptx  3%  |  .json/.ndjson/.log 7%  |  .pdf 15% (if fpdf2)

~5 % of documents include synthetic PII for redaction-pipeline testing.
Golden documents are always .txt for reliable ingestion and retrieval.
"""

from __future__ import annotations

import argparse
import csv
import json
import os
import random
import shutil
import textwrap
import uuid
from datetime import date, datetime, timedelta, timezone
from io import BytesIO
from pathlib import Path
from typing import Callable

# Optional heavy deps — degrade gracefully
try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from openpyxl import Workbook
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    from fpdf import FPDF
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

# ─────────────────────────────────────────────────────────────────────────────
# Name / data pools (deterministic, no Faker needed)
# ─────────────────────────────────────────────────────────────────────────────
FIRST_NAMES = [
    "James", "Maria", "Robert", "Patricia", "Michael", "Jennifer", "David",
    "Linda", "William", "Elizabeth", "Richard", "Barbara", "Thomas", "Susan",
    "Charles", "Jessica", "Daniel", "Sarah", "Matthew", "Karen", "Anthony",
    "Nancy", "Andrew", "Lisa", "Joshua", "Betty", "Kenneth", "Margaret",
    "Steven", "Sandra", "Kevin", "Ashley", "Brian", "Dorothy", "George",
    "Kimberly", "Timothy", "Emily", "Ronald", "Donna", "Jason", "Michelle",
    "Jeffrey", "Carol", "Ryan", "Amanda", "Jacob", "Melissa", "Gary", "Deborah",
]

LAST_NAMES = [
    "Smith", "Johnson", "Williams", "Brown", "Jones", "Garcia", "Miller",
    "Davis", "Rodriguez", "Martinez", "Hernandez", "Lopez", "Gonzalez",
    "Wilson", "Anderson", "Thomas", "Taylor", "Moore", "Jackson", "Martin",
    "Lee", "Perez", "Thompson", "White", "Harris", "Sanchez", "Clark",
    "Ramirez", "Lewis", "Robinson", "Walker", "Young", "Allen", "King",
    "Wright", "Scott", "Torres", "Nguyen", "Hill", "Flores", "Green",
    "Adams", "Nelson", "Baker", "Hall", "Rivera", "Campbell", "Mitchell",
    "Carter", "Roberts",
]

CITIES = [
    "Washington DC", "Fort Meade", "San Antonio", "Colorado Springs",
    "Norfolk", "San Diego", "Tampa", "Honolulu", "Bethesda", "Atlanta",
    "Chicago", "Houston", "Phoenix", "Seattle", "Boston", "Denver",
    "Minneapolis", "Dallas", "Portland", "Raleigh",
]

ORGS_GOV = [
    "Joint Task Force Bravo", "Cyber Command", "DISA", "NSA",
    "DIA", "SOCOM", "CENTCOM", "INDOPACOM", "EUCOM", "AFRICOM",
    "Army Futures Command", "Naval Information Warfare Center",
    "Air Force Research Laboratory", "Defense Logistics Agency",
    "Missile Defense Agency",
]

ORGS_MED = [
    "Johns Hopkins Medical Center", "Mayo Clinic", "Cleveland Clinic",
    "Massachusetts General Hospital", "Stanford Health Care",
    "Northwestern Memorial Hospital", "Cedars-Sinai Medical Center",
    "Duke University Hospital", "UCSF Medical Center",
    "Mount Sinai Hospital", "NYU Langone Health",
    "Brigham and Women's Hospital", "UCLA Medical Center",
    "Vanderbilt University Medical Center", "Emory University Hospital",
]

ORGS_ENT = [
    "Meridian Technologies", "Atlas Global Solutions", "Pinnacle Systems Inc",
    "Vanguard Digital", "Nexus Consulting Group", "Summit Analytics",
    "Horizon Partners LLC", "Apex Strategy Group", "Catalyst Innovations",
    "Sterling Operations", "Beacon Advisory", "Crestline Capital",
    "Ironbridge Solutions", "Silverline Consulting", "Keystone Ventures",
]

CLASSIFICATIONS_GOV = ["UNCLASSIFIED", "CUI", "CONFIDENTIAL", "SECRET"]
CLASSIFICATIONS_MED = ["INTERNAL - HIPAA", "INTERNAL - CLINICAL", "INTERNAL - PHI RESTRICTED"]
CLASSIFICATIONS_ENT = ["INTERNAL", "CONFIDENTIAL - BUSINESS", "INTERNAL - FINANCIAL DATA"]


def rand_name() -> str:
    return f"{random.choice(FIRST_NAMES)} {random.choice(LAST_NAMES)}"

def rand_date(days_back: int = 1000) -> str:
    d = date.today() - timedelta(days=random.randint(0, days_back))
    return d.isoformat()

def rand_money(low: int = 50_000, high: int = 50_000_000) -> str:
    v = random.randint(low, high)
    if v >= 1_000_000:
        return f"${v / 1_000_000:.1f}M"
    return f"${v:,}"

def rand_pct(low: int = 1, high: int = 99) -> str:
    return f"{random.randint(low, high)}%"

def rand_id(prefix: str = "DOC") -> str:
    return f"{prefix}-{uuid.uuid4().hex[:8].upper()}"

def rand_quarter() -> str:
    return f"Q{random.randint(1, 4)} {random.randint(2023, 2026)}"

def rand_fiscal_year() -> str:
    return f"FY{random.randint(2023, 2026)}"


# ─────────────────────────────────────────────────────────────────────────────
# PII injection (~5 % of docs)
# ─────────────────────────────────────────────────────────────────────────────
FAKE_SSNS = [
    "123-45-6789", "987-65-4321", "456-78-9012", "234-56-7890",
    "345-67-8901", "567-89-0123", "678-90-1234", "789-01-2345",
]

FAKE_EMAILS = [
    "john.doe@example.com", "jane.smith@testcorp.org", "admin@hospital-net.com",
    "mjohnson@defense.mil", "analyst@sentinel-demo.net", "cfo@acme-corp.com",
]

FAKE_PHONES = [
    "(555) 123-4567", "(555) 987-6543", "(555) 246-8135", "(555) 369-2580",
]

FAKE_CREDIT_CARDS = [
    "4111-1111-1111-1111", "5500-0000-0000-0004", "3400-000000-00009",
]

FAKE_MRN = [
    "MRN-2024-00147", "MRN-2025-03891", "MRN-2024-11562", "MRN-2025-07234",
]


def inject_pii(text: str, sector: str) -> str:
    """Inject PII markers into ~5 % of documents for redaction testing."""
    if random.random() > 0.05:
        return text

    pii_block = "\n\n--- CONTACT INFORMATION ---\n"
    name = rand_name()
    pii_block += f"Point of Contact: {name}\n"
    pii_block += f"SSN: {random.choice(FAKE_SSNS)}\n"
    pii_block += f"Email: {random.choice(FAKE_EMAILS)}\n"
    pii_block += f"Phone: {random.choice(FAKE_PHONES)}\n"

    if sector == "MEDICAL":
        pii_block += f"Medical Record Number: {random.choice(FAKE_MRN)}\n"
        pii_block += f"Patient DOB: {rand_date(days_back=20000)}\n"

    if random.random() < 0.3:
        pii_block += f"Credit Card: {random.choice(FAKE_CREDIT_CARDS)}\n"

    pii_block += "--- END CONTACT INFORMATION ---\n"

    # Insert PII block after first section break
    parts = text.split("\n\n", 2)
    if len(parts) >= 3:
        return parts[0] + "\n\n" + parts[1] + pii_block + "\n\n" + parts[2]
    return text + pii_block


# ─────────────────────────────────────────────────────────────────────────────
# GOVERNMENT templates
# ─────────────────────────────────────────────────────────────────────────────
GOV_TEMPLATES: list[Callable[[], tuple[str, str, str]]] = []

def _gov_after_action():
    op = random.choice(["IRON RESOLVE", "DARK STORM", "NORTHERN WATCH", "PACIFIC SHIELD",
                         "DESERT FALCON", "ARCTIC SENTINEL", "OCEAN GUARDIAN", "MOUNTAIN THUNDER"])
    unit = random.choice(ORGS_GOV)
    commander = rand_name()
    loc = random.choice(CITIES)
    casualties = random.randint(0, 3)
    duration = random.randint(3, 45)
    return ("after_action_report", f"After Action Report -- Operation {op}", textwrap.dedent(f"""\
        AFTER ACTION REPORT
        Classification: {random.choice(CLASSIFICATIONS_GOV)}
        Operation: {op}
        Unit: {unit}
        Commander: {commander}
        Location: {loc}
        Date: {rand_date()}
        Duration: {duration} days

        EXECUTIVE SUMMARY
        Operation {op} was conducted by {unit} under the command of {commander} in the {loc} area of operations. The operation spanned {duration} days and achieved its primary objectives with {casualties} personnel casualties.

        MISSION OBJECTIVES
        1. Establish forward operating presence in the designated area of responsibility.
        2. Conduct intelligence, surveillance, and reconnaissance operations to identify adversary positions.
        3. Coordinate with allied forces to maintain freedom of maneuver across the operational domain.
        4. Provide humanitarian assistance to affected civilian populations as directed by higher headquarters.

        TIMELINE OF EVENTS
        Day 1-{duration // 3}: Initial deployment and staging. Forces moved from marshaling areas to forward positions. Communications were established with higher headquarters and adjacent units. Intelligence preparation of the battlespace was completed.

        Day {duration // 3 + 1}-{2 * duration // 3}: Main effort execution phase. {unit} conducted coordinated operations resulting in the neutralization of {random.randint(3, 15)} threat positions. Joint fires were integrated effectively with ground maneuver elements.

        Day {2 * duration // 3 + 1}-{duration}: Consolidation and transition. Area security was established and transition to follow-on forces initiated. After-action data collection was conducted at all echelons.

        LESSONS LEARNED
        - Communications interoperability between joint elements requires pre-deployment validation and testing.
        - Logistics sustainment planning must account for extended lines of communication in austere environments.
        - Intelligence fusion at the tactical level proved critical for timely decision-making during the execution phase.
        - Civilian interaction protocols should be reinforced during pre-mission training for all personnel.

        RECOMMENDATIONS
        1. Increase allocation of ISR assets for future operations of this scope and complexity.
        2. Conduct joint interoperability exercises quarterly to maintain readiness across participating units.
        3. Update SOPs for logistics resupply in contested environments based on findings from this operation.

        APPROVED BY
        {commander}, Commanding Officer
        {unit}
        Date: {rand_date()}
    """))
GOV_TEMPLATES.append(_gov_after_action)

def _gov_intel_brief():
    threat = random.choice(["APT-29 COZY BEAR", "APT-28 FANCY BEAR", "LAZARUS GROUP",
                             "SANDWORM", "TURLA", "CHARMING KITTEN", "MUSTANG PANDA", "HAFNIUM"])
    target = random.choice(["critical infrastructure", "defense industrial base",
                             "government networks", "financial systems", "energy sector"])
    return ("intel_brief", f"Intelligence Brief -- {threat} Activity", textwrap.dedent(f"""\
        INTELLIGENCE BRIEF
        Classification: {random.choice(CLASSIFICATIONS_GOV)}
        Subject: {threat} -- Threat Assessment
        Date: {rand_date()}
        Prepared by: {rand_name()}, Intelligence Analyst
        Organization: {random.choice(ORGS_GOV)}

        THREAT OVERVIEW
        {threat} has been observed conducting sustained cyber operations targeting {target} across multiple allied nations. Activity has increased {rand_pct(15, 85)} over the past reporting period. This assessment reflects analysis of {random.randint(50, 500)} unique indicators of compromise collected across partner intelligence networks.

        INDICATORS OF COMPROMISE
        - Command and control infrastructure observed on {random.randint(5, 30)} unique IP ranges spanning {random.randint(3, 12)} countries.
        - Spear-phishing campaigns leveraging credential harvesting techniques against {random.randint(100, 2000)} targeted accounts.
        - Use of zero-day vulnerabilities in widely deployed enterprise software products, with {random.randint(2, 8)} confirmed exploits.
        - Lateral movement using living-off-the-land techniques to evade endpoint detection and response tools.

        ASSESSMENT
        {threat} is assessed with HIGH CONFIDENCE to be state-sponsored with tasking focused on intelligence collection against {target}. The group maintains sophisticated operational security and routinely rotates infrastructure to complicate attribution. Current intelligence suggests preparation for escalated operations in the next {random.randint(30, 180)} days.

        RECOMMENDED COUNTERMEASURES
        1. Implement enhanced monitoring on all network egress points for known {threat} indicators.
        2. Conduct targeted threat hunting operations focused on lateral movement patterns.
        3. Accelerate patching cycles for identified vulnerability categories.
        4. Increase information sharing with sector partners and allied intelligence services.

        DISTRIBUTION: {random.choice(CLASSIFICATIONS_GOV)} // REL TO FVEY
        Prepared by: {rand_name()}
    """))
GOV_TEMPLATES.append(_gov_intel_brief)

def _gov_policy_memo():
    subject = random.choice([
        "Zero Trust Architecture Implementation", "Cloud Migration Security Requirements",
        "Insider Threat Program Enhancements", "Supply Chain Risk Management",
        "Continuous Authorization to Operate", "Cross-Domain Solution Standards",
        "Mobile Device Management Policy Update", "Data Loss Prevention Requirements",
    ])
    return ("policy_memo", f"Policy Memorandum -- {subject}", textwrap.dedent(f"""\
        MEMORANDUM FOR: All {random.choice(ORGS_GOV)} Personnel
        FROM: {rand_name()}, Director of Policy
        SUBJECT: {subject}
        DATE: {rand_date()}
        Classification: {random.choice(CLASSIFICATIONS_GOV)}

        1. PURPOSE
        This memorandum establishes updated policy guidance for {subject.lower()} across all organizational components. This policy supersedes previous guidance dated {rand_date(days_back=500)} and is effective immediately upon distribution.

        2. BACKGROUND
        Recent assessments have identified the need to modernize our approach to {subject.lower()}. Industry best practices and evolving threat landscapes require a comprehensive update to existing frameworks. A review of {random.randint(10, 50)} related incidents over the past {random.randint(6, 24)} months informed the development of this policy.

        3. POLICY
        a. All systems classified at the {random.choice(CLASSIFICATIONS_GOV)} level and above shall implement the requirements outlined in Enclosure 1 within {random.randint(90, 365)} days of this memorandum.
        b. System owners are responsible for developing implementation plans and submitting them to the Chief Information Security Officer within 30 days.
        c. Annual compliance assessments shall be conducted by independent evaluators beginning in {rand_fiscal_year()}.
        d. Exceptions to this policy require written approval from the authorizing official with documentation of compensating controls and risk acceptance.

        4. RESPONSIBILITIES
        a. Chief Information Officer: Overall policy oversight and implementation guidance.
        b. Chief Information Security Officer: Compliance monitoring and assessment coordination.
        c. System Owners: Implementation planning, execution, and continuous monitoring.
        d. All Personnel: Compliance with updated procedures and completion of required training.

        5. EFFECTIVE DATE
        This policy is effective upon signature and remains in effect until superseded or rescinded.

        {rand_name()}
        Director of Policy
        {random.choice(ORGS_GOV)}
    """))
GOV_TEMPLATES.append(_gov_policy_memo)

def _gov_procurement():
    item = random.choice([
        "Next-Generation Firewall Systems", "Endpoint Detection and Response Platform",
        "Secure Cloud Hosting Services", "Cryptographic Key Management System",
        "Identity and Access Management Suite", "Network Traffic Analysis Tools",
        "Vulnerability Assessment Platform", "Security Information and Event Management",
    ])
    vendor = random.choice(["Lockheed Martin", "Raytheon", "Northrop Grumman",
                             "General Dynamics", "BAE Systems", "L3Harris", "Leidos", "SAIC"])
    cost = rand_money(500_000, 25_000_000)
    return ("procurement_summary", f"Procurement Summary -- {item}", textwrap.dedent(f"""\
        PROCUREMENT SUMMARY
        Classification: {random.choice(CLASSIFICATIONS_GOV)}
        Document ID: {rand_id("PROC")}
        Date: {rand_date()}
        Prepared by: {rand_name()}, Contracting Officer

        ITEM: {item}
        VENDOR: {vendor}
        CONTRACT VALUE: {cost}
        CONTRACT TYPE: {random.choice(["Firm Fixed Price", "Cost Plus Fixed Fee", "Time and Materials", "Indefinite Delivery/Indefinite Quantity"])}
        PERIOD OF PERFORMANCE: {rand_fiscal_year()} through {rand_fiscal_year()}

        JUSTIFICATION
        The acquisition of {item} is necessary to address identified capability gaps in the organization's cyber defense posture. A market survey of {random.randint(4, 12)} vendors was conducted, and {vendor} was selected based on technical merit, past performance, and cost reasonableness.

        TECHNICAL EVALUATION
        {vendor}'s solution demonstrated superior performance in {random.randint(3, 7)} of {random.randint(8, 12)} evaluated criteria. Key strengths included integration with existing infrastructure, scalability to meet projected growth, and compliance with applicable security standards. The solution met all mandatory requirements and exceeded expectations in {random.randint(2, 5)} optional areas.

        BUDGET IMPACT
        Funding is available from the {rand_fiscal_year()} operations and maintenance account. The total lifecycle cost over the {random.randint(3, 7)}-year period of performance is estimated at {rand_money(1_000_000, 50_000_000)}, including maintenance, training, and sustainment.

        RISK ASSESSMENT
        Overall procurement risk is assessed as {random.choice(["LOW", "MODERATE", "LOW-MODERATE"])}. Primary risk factors include supply chain dependencies and integration timeline constraints.

        APPROVED BY
        {rand_name()}, Contracting Officer
        {rand_name()}, Program Manager
    """))
GOV_TEMPLATES.append(_gov_procurement)

def _gov_cyber_incident():
    incident_type = random.choice(["Ransomware Attack", "Data Exfiltration", "Phishing Compromise",
                                    "Insider Threat", "DDoS Attack", "Supply Chain Compromise",
                                    "Credential Stuffing", "Zero-Day Exploitation"])
    severity = random.choice(["CRITICAL", "HIGH", "MODERATE"])
    systems = random.randint(1, 150)
    return ("cyber_incident_report", f"Cyber Incident Report -- {incident_type}", textwrap.dedent(f"""\
        CYBER INCIDENT REPORT
        Classification: {random.choice(CLASSIFICATIONS_GOV)}
        Incident ID: {rand_id("INC")}
        Date Detected: {rand_date(days_back=90)}
        Severity: {severity}
        Incident Type: {incident_type}
        Reporting Organization: {random.choice(ORGS_GOV)}

        INCIDENT SUMMARY
        A {incident_type.lower()} incident was detected on {rand_date(days_back=90)} affecting approximately {systems} systems within the organization's network. The incident was initially identified through automated monitoring systems and escalated to the incident response team within {random.randint(5, 120)} minutes of detection.

        AFFECTED SYSTEMS
        - {systems} endpoints confirmed affected across {random.randint(1, 8)} network segments.
        - {random.randint(0, 10)} servers in the production environment showed indicators of compromise.
        - {random.randint(0, 5)} user accounts were confirmed compromised and immediately disabled.
        - Network segments {random.choice(["A", "B", "C"])}-{random.randint(1, 20)} through {random.choice(["A", "B", "C"])}-{random.randint(21, 40)} were isolated during containment.

        RESPONSE ACTIONS
        1. CONTAIN: Affected network segments isolated within {random.randint(15, 180)} minutes. All compromised accounts disabled. Firewall rules updated to block identified command and control infrastructure.
        2. ERADICATE: Malicious artifacts removed from {systems} systems. Full reimaging conducted on {random.randint(1, systems)} critically affected endpoints.
        3. RECOVER: Systems restored from verified clean backups. Enhanced monitoring deployed across all affected segments. User credentials reset for {random.randint(50, 500)} accounts as a precautionary measure.

        LESSONS LEARNED
        - Detection time could be improved with enhanced behavioral analytics at the endpoint level.
        - Incident response playbooks should be updated to include procedures for this specific attack vector.
        - Segmentation controls proved effective in limiting lateral movement during the initial containment phase.

        STATUS: {random.choice(["RESOLVED", "MONITORING", "UNDER INVESTIGATION"])}
        Incident Commander: {rand_name()}
    """))
GOV_TEMPLATES.append(_gov_cyber_incident)

def _gov_risk_assessment():
    system = random.choice(["Joint Warfighting Cloud Capability", "Defense Enterprise Email",
                             "Global Command and Control System", "Tactical Data Links Network",
                             "Personnel Security System", "Logistics Information System"])
    return ("risk_assessment", f"Risk Assessment -- {system}", textwrap.dedent(f"""\
        RISK ASSESSMENT
        Classification: {random.choice(CLASSIFICATIONS_GOV)}
        System: {system}
        Assessment Date: {rand_date()}
        Assessor: {rand_name()}, ISSM
        Organization: {random.choice(ORGS_GOV)}

        EXECUTIVE SUMMARY
        This risk assessment evaluates the security posture of {system} in accordance with NIST SP 800-30 and DoD RMF requirements. The overall risk level is assessed as {random.choice(["LOW", "MODERATE", "HIGH"])} based on analysis of {random.randint(15, 80)} identified vulnerabilities and {random.randint(5, 25)} threat scenarios.

        SYSTEM DESCRIPTION
        {system} is a {random.choice(["mission-critical", "mission-essential", "supporting"])} system that provides {random.choice(["command and control", "communications", "logistics", "intelligence"])} capabilities to {random.randint(500, 50000)} users across {random.randint(3, 50)} sites.

        RISK FINDINGS
        HIGH RISK ({random.randint(1, 5)} findings):
        - Unpatched vulnerabilities in core application components exceed the {random.randint(15, 45)}-day remediation window.
        - Privileged access management controls do not meet current STIG requirements.

        MODERATE RISK ({random.randint(3, 12)} findings):
        - Audit logging coverage is incomplete for {rand_pct(5, 30)} of system transactions.
        - Backup and recovery procedures have not been tested within the required {random.randint(90, 365)}-day cycle.
        - Multi-factor authentication is not enforced for all remote access paths.

        LOW RISK ({random.randint(5, 20)} findings):
        - Documentation gaps in system security plans and procedures.
        - Minor configuration deviations from established baselines.

        RECOMMENDATIONS
        1. Immediately remediate all HIGH risk findings within {random.randint(14, 30)} days.
        2. Develop plan of action and milestones for all MODERATE findings with completion within {random.randint(60, 180)} days.
        3. Address LOW risk findings during the next scheduled maintenance window.

        AUTHORIZATION RECOMMENDATION
        Based on this assessment, {system} is recommended for {random.choice(["continued authorization with conditions", "authorization to operate", "denial of authorization pending remediation"])}.

        {rand_name()}, Authorizing Official
    """))
GOV_TEMPLATES.append(_gov_risk_assessment)


# ─────────────────────────────────────────────────────────────────────────────
# MEDICAL templates
# ─────────────────────────────────────────────────────────────────────────────
MED_TEMPLATES: list[Callable[[], tuple[str, str, str]]] = []

def _med_clinical_trial():
    drug = random.choice(["MED-7742", "BIO-3391", "THR-2210", "IMM-5568",
                           "ONC-9913", "NEU-4427", "CAR-6651", "END-1189"])
    phase = random.choice(["Phase I", "Phase II", "Phase III", "Phase IV"])
    condition = random.choice(["Type 2 Diabetes", "Non-Small Cell Lung Cancer",
                                "Rheumatoid Arthritis", "Major Depressive Disorder",
                                "Chronic Heart Failure", "Alzheimer's Disease",
                                "Multiple Sclerosis", "Crohn's Disease"])
    enrolled = random.randint(50, 5000)
    return ("clinical_trial_summary", f"Clinical Trial Summary -- {drug} {phase}", textwrap.dedent(f"""\
        CLINICAL TRIAL SUMMARY
        Classification: {random.choice(CLASSIFICATIONS_MED)}
        Trial ID: {rand_id("CT")}
        Protocol: {drug}-{phase.replace(' ', '')}-{random.randint(100, 999)}
        Date: {rand_date()}
        Principal Investigator: Dr. {rand_name()}
        Institution: {random.choice(ORGS_MED)}

        STUDY OVERVIEW
        This {phase} clinical trial evaluates the safety and efficacy of {drug} for the treatment of {condition}. The study enrolled {enrolled} participants across {random.randint(5, 50)} clinical sites over a {random.randint(6, 36)}-month period. Enrollment criteria included adults aged {random.randint(18, 40)}-{random.randint(65, 85)} with confirmed diagnosis of {condition} and no contraindicated comorbidities.

        STUDY DESIGN
        - Design: {random.choice(["Randomized, double-blind, placebo-controlled", "Open-label, single-arm", "Randomized, active-comparator"])}
        - Primary Endpoint: {random.choice(["Overall survival at 12 months", "Reduction in symptom severity score", "Disease progression-free interval", "Change in biomarker levels from baseline"])}
        - Treatment Arms: {drug} ({random.choice(["50mg", "100mg", "200mg", "5mg/kg"])}) vs. {random.choice(["placebo", "standard of care", "active comparator"])}
        - Duration: {random.randint(12, 52)} weeks with {random.randint(4, 24)}-week follow-up

        EFFICACY RESULTS
        The primary endpoint was {random.choice(["met", "met with statistical significance", "not met"])} (p={random.choice(["<0.001", "<0.01", "<0.05", "0.07", "0.12"])}). The treatment group showed a {rand_pct(5, 45)} improvement in the primary outcome measure compared to the control group. Secondary endpoints demonstrated consistent trends favoring the {drug} treatment arm in {random.randint(2, 5)} of {random.randint(4, 8)} measures evaluated.

        SAFETY PROFILE
        - Adverse Events: {random.randint(10, 60)}% of participants reported at least one adverse event.
        - Serious Adverse Events: {random.randint(1, 15)}% incidence, with {random.randint(0, 5)} events considered possibly related to study drug.
        - Most common AEs: {random.choice(["nausea", "headache", "fatigue"])}, {random.choice(["dizziness", "insomnia", "rash"])}, {random.choice(["arthralgia", "diarrhea", "cough"])}
        - Discontinuation due to AEs: {random.randint(2, 20)}%

        CONCLUSIONS
        {drug} demonstrated a {random.choice(["favorable", "acceptable", "manageable"])} safety profile in this {phase} study. {random.choice(["Results support advancement to the next phase of clinical development.", "Further investigation is warranted to optimize dosing and identify responsive patient subgroups.", "The data support submission of a supplemental application to regulatory authorities."])}

        HIPAA NOTICE: This document contains protected health information in aggregate form.
        PI Signature: Dr. {rand_name()}
    """))
MED_TEMPLATES.append(_med_clinical_trial)

def _med_care_protocol():
    condition = random.choice(["Acute Myocardial Infarction", "Sepsis",
                                "Diabetic Ketoacidosis", "Acute Stroke",
                                "Pneumonia", "Hip Fracture", "Heart Failure Exacerbation",
                                "Acute Kidney Injury"])
    return ("care_protocol", f"Care Protocol -- {condition}", textwrap.dedent(f"""\
        CLINICAL CARE PROTOCOL
        Classification: {random.choice(CLASSIFICATIONS_MED)}
        Protocol ID: {rand_id("PROT")}
        Condition: {condition}
        Effective Date: {rand_date()}
        Approved by: Dr. {rand_name()}, Chief Medical Officer
        Institution: {random.choice(ORGS_MED)}
        Review Cycle: Annual

        PURPOSE
        This protocol establishes standardized clinical pathways for the management of {condition} at {random.choice(ORGS_MED)}. It is intended to reduce variability in care delivery, improve patient outcomes, and ensure compliance with current evidence-based guidelines.

        INITIAL ASSESSMENT
        1. Obtain vital signs including blood pressure, heart rate, respiratory rate, oxygen saturation, and temperature within {random.randint(5, 15)} minutes of presentation.
        2. Perform targeted physical examination with focus on {random.choice(["cardiovascular", "respiratory", "neurological", "abdominal"])} systems.
        3. Order initial laboratory panel: CBC, BMP, troponin, BNP, lactate, coagulation studies as clinically indicated.
        4. Obtain {random.choice(["12-lead ECG", "chest radiograph", "CT scan", "point-of-care ultrasound"])} within {random.randint(15, 60)} minutes.

        TREATMENT PATHWAY
        Hour 0-{random.randint(1, 4)}:
        - Establish IV access with {random.choice(["normal saline", "lactated Ringer's"])} at {random.randint(100, 500)} mL/hr.
        - Administer initial medications per standing orders.
        - Consult specialty services as indicated by clinical presentation.
        - Communicate with patient and family regarding diagnosis and treatment plan.

        Hour {random.randint(4, 8)}-{random.randint(12, 24)}:
        - Reassess clinical status every {random.randint(2, 4)} hours.
        - Titrate medications based on clinical response and laboratory results.
        - Initiate multidisciplinary care coordination including pharmacy, nursing, social work, and rehabilitation services.

        DISCHARGE CRITERIA
        - Clinical stability for minimum {random.randint(12, 48)} hours.
        - Oral medication tolerance demonstrated.
        - Follow-up appointment scheduled within {random.randint(3, 14)} days.
        - Patient education completed with written discharge instructions provided.
        - Home medication reconciliation performed by clinical pharmacist.

        QUALITY METRICS
        - Door-to-treatment time target: <{random.randint(30, 120)} minutes
        - 30-day readmission rate target: <{rand_pct(5, 15)}
        - Patient satisfaction score target: >{rand_pct(80, 95)}

        HIPAA NOTICE: All patient information must be handled per institutional PHI policies.
    """))
MED_TEMPLATES.append(_med_care_protocol)

def _med_quality_metrics():
    dept = random.choice(["Emergency Department", "Intensive Care Unit", "Surgical Services",
                           "Pediatrics", "Oncology", "Cardiology", "Orthopedics", "Obstetrics"])
    quarter = rand_quarter()
    return ("quality_metrics", f"Quality Metrics Report -- {dept} {quarter}", textwrap.dedent(f"""\
        QUALITY METRICS REPORT
        Classification: {random.choice(CLASSIFICATIONS_MED)}
        Department: {dept}
        Reporting Period: {quarter}
        Prepared by: {rand_name()}, Quality Improvement Coordinator
        Institution: {random.choice(ORGS_MED)}

        EXECUTIVE SUMMARY
        The {dept} achieved {random.randint(3, 8)} of {random.randint(8, 12)} quality benchmarks during {quarter}. Overall performance improved {rand_pct(1, 15)} compared to the prior quarter. Patient safety incidents decreased by {rand_pct(5, 30)}, and staff compliance with hand hygiene protocols reached {rand_pct(85, 99)}.

        KEY PERFORMANCE INDICATORS
        1. Patient Satisfaction (HCAHPS): {rand_pct(70, 98)} (Target: 85%)
        2. Average Length of Stay: {random.uniform(2.0, 12.0):.1f} days (Target: {random.uniform(2.0, 8.0):.1f} days)
        3. 30-Day Readmission Rate: {rand_pct(3, 18)} (Target: <10%)
        4. Hospital-Acquired Infection Rate: {random.uniform(0.1, 3.5):.1f} per 1,000 patient days
        5. Medication Error Rate: {random.uniform(0.01, 1.5):.2f} per 1,000 doses administered
        6. Fall Rate: {random.uniform(0.5, 4.0):.1f} per 1,000 patient days
        7. Door-to-Provider Time: {random.randint(8, 45)} minutes (Target: <20 minutes)
        8. Staff Turnover Rate: {rand_pct(5, 25)} annualized

        IMPROVEMENT INITIATIVES
        - Implemented bedside shift reporting in {random.randint(3, 8)} units, resulting in {rand_pct(10, 35)} reduction in patient-reported communication gaps.
        - Deployed real-time location system for equipment tracking, reducing search time by {rand_pct(20, 50)}.
        - Launched multidisciplinary rounds with pharmacy participation, decreasing medication reconciliation errors by {rand_pct(15, 40)}.

        ACTION ITEMS
        1. Address staffing shortages in {dept} through targeted recruitment -- deadline {rand_date(days_back=0)}.
        2. Implement automated sepsis screening in EHR -- target completion {rand_date(days_back=0)}.
        3. Conduct root cause analysis for readmission outliers identified this quarter.

        HIPAA NOTICE: Aggregate data only. No individual patient data included.
    """))
MED_TEMPLATES.append(_med_quality_metrics)

def _med_safety_incident():
    incident = random.choice(["Medication Error", "Patient Fall", "Surgical Site Infection",
                               "Wrong-Site Procedure Near Miss", "Blood Transfusion Reaction",
                               "Equipment Malfunction", "Diagnostic Delay", "Communication Failure"])
    return ("safety_incident", f"Patient Safety Incident Report -- {incident}", textwrap.dedent(f"""\
        PATIENT SAFETY INCIDENT REPORT
        Classification: {random.choice(CLASSIFICATIONS_MED)}
        Incident ID: {rand_id("PSI")}
        Date of Incident: {rand_date(days_back=180)}
        Severity: {random.choice(["Near Miss", "No Harm", "Temporary Harm", "Permanent Harm"])}
        Incident Type: {incident}
        Department: {random.choice(["Emergency", "ICU", "Med-Surg", "OR", "Pharmacy", "Radiology"])}
        Institution: {random.choice(ORGS_MED)}

        INCIDENT DESCRIPTION
        A {incident.lower()} event occurred on {rand_date(days_back=180)} at approximately {random.randint(1, 12)}:{random.randint(10, 59)} {random.choice(["AM", "PM"])}. The incident was identified by {random.choice(["nursing staff", "attending physician", "pharmacist", "patient", "automated alert system"])} and immediately reported through the institutional safety reporting system.

        ROOT CAUSE ANALYSIS
        Contributing factors identified through investigation include:
        - {random.choice(["Process/System failure", "Communication breakdown", "Environmental factor", "Training gap"])} as the primary contributing factor.
        - {random.choice(["Staffing levels", "Workload", "Fatigue", "Distraction"])} identified as a secondary factor.
        - Existing safeguards {random.choice(["functioned as designed but were insufficient", "were bypassed due to workflow constraints", "were not applicable to this scenario"])}.

        IMMEDIATE ACTIONS TAKEN
        1. Patient assessed and stabilized. {random.choice(["No additional treatment required.", "Appropriate treatment administered.", "Specialist consultation obtained."])}
        2. Incident reported to attending physician, charge nurse, and department leadership.
        3. Equipment and medications involved quarantined for investigation.
        4. Patient and family notified per disclosure policy.

        CORRECTIVE ACTIONS
        1. {random.choice(["Updated medication administration protocol", "Enhanced fall prevention bundle", "Revised surgical safety checklist", "Improved handoff communication template"])} -- Implementation by {rand_date(days_back=0)}.
        2. Staff re-education completed for all affected personnel within {random.randint(7, 30)} days.
        3. System modification to add {random.choice(["additional verification step", "automated alert", "forcing function", "independent double-check"])} -- Target completion {rand_date(days_back=0)}.

        HIPAA NOTICE: All identifiers have been removed. This report contains de-identified data only.
        Report filed by: {rand_name()}, Risk Management
    """))
MED_TEMPLATES.append(_med_safety_incident)

def _med_formulary():
    committee_date = rand_date()
    additions = random.randint(3, 12)
    removals = random.randint(1, 5)
    return ("formulary_update", f"Pharmacy & Therapeutics Formulary Update", textwrap.dedent(f"""\
        PHARMACY & THERAPEUTICS COMMITTEE
        FORMULARY UPDATE
        Classification: {random.choice(CLASSIFICATIONS_MED)}
        Committee Meeting Date: {committee_date}
        Institution: {random.choice(ORGS_MED)}
        Chair: Dr. {rand_name()}, PharmD

        SUMMARY OF ACTIONS
        The P&T Committee reviewed {random.randint(10, 40)} drug monographs and took the following actions:
        - {additions} medications added to formulary
        - {removals} medications removed from formulary
        - {random.randint(2, 8)} therapeutic substitution protocols updated
        - {random.randint(1, 5)} clinical guidelines revised

        ADDITIONS TO FORMULARY
        1. {random.choice(["Semaglutide", "Tirzepatide", "Pembrolizumab", "Osimertinib"])} -- Approved for {random.choice(["Type 2 Diabetes", "obesity management", "non-small cell lung cancer", "chronic kidney disease"])}. Estimated annual cost impact: {rand_money(100_000, 5_000_000)}.
        2. {random.choice(["Tofacitinib", "Baricitinib", "Upadacitinib"])} -- Restricted use for {random.choice(["rheumatoid arthritis", "atopic dermatitis", "ulcerative colitis"])} after failure of first-line therapy.
        3. {random.choice(["Empagliflozin", "Dapagliflozin", "Canagliflozin"])} -- Added for heart failure indication regardless of diabetes status.

        COST ANALYSIS
        Total formulary cost impact for the upcoming fiscal year is projected at {rand_money(500_000, 10_000_000)}. The committee approved {random.randint(2, 6)} therapeutic interchange protocols expected to generate savings of {rand_money(100_000, 3_000_000)} annually.

        NEXT MEETING: {rand_date(days_back=0)}
        Minutes prepared by: {rand_name()}, PharmD, Clinical Pharmacy Specialist
    """))
MED_TEMPLATES.append(_med_formulary)


# ─────────────────────────────────────────────────────────────────────────────
# ENTERPRISE templates
# ─────────────────────────────────────────────────────────────────────────────
ENT_TEMPLATES: list[Callable[[], tuple[str, str, str]]] = []

def _ent_strategy_memo():
    initiative = random.choice([
        "Digital Transformation Program", "Cloud-First Migration Strategy",
        "Customer Experience Modernization", "AI-Powered Operations Initiative",
        "Global Expansion Framework", "Sustainability and ESG Program",
        "Workforce of the Future Initiative", "Data Monetization Strategy",
    ])
    return ("strategy_memo", f"Strategic Initiative -- {initiative}", textwrap.dedent(f"""\
        STRATEGIC INITIATIVE MEMO
        Classification: {random.choice(CLASSIFICATIONS_ENT)}
        Initiative: {initiative}
        Date: {rand_date()}
        Author: {rand_name()}, Chief Strategy Officer
        Organization: {random.choice(ORGS_ENT)}

        EXECUTIVE SUMMARY
        The {initiative} represents a {rand_money(1_000_000, 50_000_000)} investment over {random.randint(2, 5)} years to fundamentally transform our competitive position. This initiative is projected to deliver {rand_pct(15, 45)} improvement in operational efficiency and generate {rand_money(5_000_000, 100_000_000)} in incremental revenue by {rand_fiscal_year()}.

        STRATEGIC RATIONALE
        Market analysis indicates that {rand_pct(40, 80)} of industry peers have already initiated similar programs. Our current trajectory places us at risk of falling behind in key capability areas. Customer surveys indicate {rand_pct(50, 85)} of enterprise clients consider these capabilities a prerequisite for contract renewal.

        KEY WORKSTREAMS
        1. Technology Foundation ({rand_money(500_000, 10_000_000)}): Modernize core infrastructure to support {random.choice(["cloud-native", "AI-driven", "data-centric", "platform-based"])} operating model.
        2. Process Transformation ({rand_money(300_000, 5_000_000)}): Redesign {random.randint(5, 20)} core business processes to eliminate manual steps and reduce cycle time by {rand_pct(30, 60)}.
        3. People and Culture ({rand_money(200_000, 3_000_000)}): Upskill {random.randint(100, 2000)} employees and recruit {random.randint(10, 100)} specialized roles to support new capabilities.
        4. Change Management ({rand_money(100_000, 2_000_000)}): Execute comprehensive communications and adoption program across {random.randint(3, 15)} business units.

        RISK FACTORS
        - Execution complexity across {random.randint(3, 15)} business units and {random.randint(5, 30)} geographic locations.
        - Talent acquisition in a competitive market for specialized skills.
        - Technology integration with {random.randint(5, 25)} legacy systems requiring migration or retirement.

        GOVERNANCE
        Steering committee meets {random.choice(["monthly", "bi-weekly"])} with executive sponsor {rand_name()}, CEO.
        Program Director: {rand_name()}
    """))
ENT_TEMPLATES.append(_ent_strategy_memo)

def _ent_quarterly_review():
    quarter = rand_quarter()
    revenue = random.randint(10, 500)
    return ("quarterly_business_review", f"Quarterly Business Review -- {quarter}", textwrap.dedent(f"""\
        QUARTERLY BUSINESS REVIEW
        Classification: {random.choice(CLASSIFICATIONS_ENT)}
        Period: {quarter}
        Date: {rand_date()}
        Prepared by: {rand_name()}, VP Finance
        Organization: {random.choice(ORGS_ENT)}

        FINANCIAL SUMMARY
        Revenue: ${revenue}M ({random.choice(["+", "-"])}{rand_pct(1, 15)} vs. prior quarter)
        Gross Margin: {rand_pct(35, 75)}
        EBITDA: ${random.randint(1, revenue // 2)}M
        Operating Cash Flow: ${random.randint(1, revenue // 3)}M
        Headcount: {random.randint(200, 10000)} FTEs ({random.choice(["+", "-"])}{random.randint(5, 200)} vs. prior quarter)

        REVENUE BY SEGMENT
        - Professional Services: ${random.randint(1, revenue // 2)}M ({rand_pct(1, 15)} growth)
        - Software Licensing: ${random.randint(1, revenue // 3)}M ({rand_pct(1, 20)} growth)
        - Managed Services: ${random.randint(1, revenue // 4)}M ({rand_pct(1, 25)} growth)
        - Training and Advisory: ${random.randint(1, revenue // 5)}M ({rand_pct(1, 10)} growth)

        KEY WINS
        1. Signed {random.randint(1, 5)}-year contract with {random.choice(["Fortune 500 client", "federal agency", "healthcare system", "financial institution"])} valued at {rand_money(500_000, 20_000_000)}.
        2. Expanded existing relationship with {random.choice(ORGS_ENT)} adding {random.randint(2, 8)} new service lines.
        3. Launched new {random.choice(["AI analytics", "cloud security", "data platform", "automation"])} offering with {random.randint(5, 30)} early adopter clients.

        PIPELINE
        Qualified pipeline: {rand_money(10_000_000, 200_000_000)}
        Win rate (trailing 4 quarters): {rand_pct(20, 55)}
        Average deal size: {rand_money(100_000, 5_000_000)}
        Sales cycle: {random.randint(30, 180)} days average

        OUTLOOK
        We project {rand_pct(5, 20)} revenue growth in the next quarter based on pipeline conversion rates and confirmed renewals. Key risks include {random.choice(["macroeconomic headwinds", "competitive pricing pressure", "talent retention", "regulatory changes"])}.

        Reviewed by: {rand_name()}, CEO
    """))
ENT_TEMPLATES.append(_ent_quarterly_review)

def _ent_compliance_audit():
    framework = random.choice(["SOC 2 Type II", "ISO 27001", "PCI DSS", "NIST CSF",
                                 "GDPR", "CCPA", "FedRAMP", "CMMC Level 2"])
    return ("compliance_audit", f"Compliance Audit Report -- {framework}", textwrap.dedent(f"""\
        COMPLIANCE AUDIT REPORT
        Classification: {random.choice(CLASSIFICATIONS_ENT)}
        Framework: {framework}
        Audit Period: {rand_date()} through {rand_date(days_back=0)}
        Auditor: {rand_name()}, CISA, CISSP
        Organization: {random.choice(ORGS_ENT)}

        EXECUTIVE SUMMARY
        This report presents the findings of the {framework} compliance assessment conducted for {random.choice(ORGS_ENT)}. The assessment evaluated {random.randint(50, 200)} controls across {random.randint(5, 15)} control domains. Overall compliance rate is {rand_pct(75, 98)}, with {random.randint(0, 5)} critical findings and {random.randint(2, 15)} non-critical observations.

        SCOPE
        The audit covered all production systems, data processing activities, and supporting infrastructure. {random.randint(3, 12)} business applications, {random.randint(5, 30)} servers, and {random.randint(2, 8)} third-party integrations were included in the assessment scope.

        FINDINGS SUMMARY
        Compliant: {random.randint(80, 195)} controls
        Partially Compliant: {random.randint(3, 15)} controls
        Non-Compliant: {random.randint(0, 5)} controls
        Not Applicable: {random.randint(0, 10)} controls

        CRITICAL FINDINGS
        1. {random.choice(["Access management procedures require enhancement", "Data encryption at rest does not meet requirements for all classified data stores", "Incident response plan has not been tested within the required period", "Vendor risk management program requires formal documentation"])}
           Remediation deadline: {rand_date(days_back=0)}

        NON-CRITICAL OBSERVATIONS
        1. Policy documentation should be updated to reflect current operational procedures.
        2. Security awareness training completion rate is {rand_pct(70, 89)}, below the {rand_pct(90, 95)} target.
        3. Asset inventory requires reconciliation with configuration management database.

        REMEDIATION PLAN
        All findings have been assigned owners with remediation timelines ranging from {random.randint(14, 30)} to {random.randint(60, 180)} days. Progress will be tracked through the GRC platform with monthly reporting to the audit committee.

        OPINION
        Based on our assessment, {random.choice(ORGS_ENT)} has {random.choice(["achieved", "substantially achieved", "partially achieved"])} compliance with {framework} requirements.

        Lead Auditor: {rand_name()}, CISA
    """))
ENT_TEMPLATES.append(_ent_compliance_audit)

def _ent_vendor_assessment():
    vendor = random.choice(["CloudStar Hosting", "DataVault Solutions", "SecureNet Partners",
                             "NexaForce Technologies", "OmniTech Services", "PrismWare Inc",
                             "Fortex Consulting", "Luminary Digital"])
    service = random.choice(["Cloud Infrastructure", "Managed Security Services",
                              "Data Analytics Platform", "HR Information System",
                              "Customer Relationship Management", "Enterprise Resource Planning"])
    return ("vendor_assessment", f"Vendor Assessment -- {vendor}", textwrap.dedent(f"""\
        VENDOR ASSESSMENT REPORT
        Classification: {random.choice(CLASSIFICATIONS_ENT)}
        Vendor: {vendor}
        Service: {service}
        Assessment Date: {rand_date()}
        Assessor: {rand_name()}, Procurement Manager
        Organization: {random.choice(ORGS_ENT)}

        VENDOR OVERVIEW
        {vendor} provides {service.lower()} to approximately {random.randint(100, 5000)} enterprise clients globally. The company was founded in {random.randint(2005, 2020)} and has {random.randint(50, 5000)} employees across {random.randint(2, 15)} offices. Annual revenue is estimated at {rand_money(5_000_000, 500_000_000)}.

        EVALUATION CRITERIA
        1. Technical Capability: {rand_pct(60, 98)}/100
        2. Financial Stability: {rand_pct(50, 95)}/100
        3. Security Posture: {rand_pct(65, 99)}/100
        4. Customer References: {rand_pct(70, 98)}/100
        5. Cost Competitiveness: {rand_pct(55, 95)}/100
        6. SLA Commitments: {rand_pct(70, 99)}/100

        OVERALL SCORE: {rand_pct(65, 95)}/100

        STRENGTHS
        - Strong technical capabilities with demonstrated expertise in {service.lower()}.
        - {random.choice(["SOC 2 Type II", "ISO 27001", "FedRAMP"])} certified with clean audit history.
        - Customer satisfaction scores averaging {rand_pct(80, 96)} across surveyed references.

        CONCERNS
        - Financial statements show declining margins over the past {random.randint(2, 4)} quarters.
        - Key personnel turnover of {rand_pct(10, 30)} in the past year.
        - Sub-processor management practices require improvement.

        RECOMMENDATION: {random.choice(["APPROVED", "APPROVED WITH CONDITIONS", "FURTHER EVALUATION REQUIRED"])}
        Contract Value: {rand_money(50_000, 5_000_000)} annually
        Contract Term: {random.randint(1, 5)} years

        Approved by: {rand_name()}, VP Procurement
    """))
ENT_TEMPLATES.append(_ent_vendor_assessment)

def _ent_security_incident():
    incident = random.choice(["Phishing Campaign", "Data Breach Investigation",
                               "Malware Detection", "Unauthorized Access Attempt",
                               "Social Engineering Attack", "Credential Compromise"])
    return ("security_incident", f"Security Incident Report -- {incident}", textwrap.dedent(f"""\
        SECURITY INCIDENT REPORT
        Classification: {random.choice(CLASSIFICATIONS_ENT)}
        Incident ID: {rand_id("SEC")}
        Incident Type: {incident}
        Date Detected: {rand_date(days_back=90)}
        Severity: {random.choice(["P1 - Critical", "P2 - High", "P3 - Medium"])}
        Status: {random.choice(["Resolved", "Under Investigation", "Monitoring"])}
        Organization: {random.choice(ORGS_ENT)}

        INCIDENT SUMMARY
        A {incident.lower()} was detected on {rand_date(days_back=90)} by {random.choice(["SOC analysts", "automated SIEM alerts", "employee report", "endpoint detection system"])}. The incident affected {random.randint(1, 500)} user accounts and {random.randint(0, 50)} systems. Containment was achieved within {random.randint(15, 720)} minutes of initial detection.

        IMPACT ASSESSMENT
        - Data Exposure: {random.choice(["No confirmed data exposure", "Limited exposure of internal business data", "Customer data potentially accessed"])}
        - Systems Affected: {random.randint(1, 50)} endpoints, {random.randint(0, 10)} servers
        - Business Impact: {random.choice(["Minimal operational disruption", "Moderate impact to operations for 4 hours", "Significant disruption requiring BCP activation"])}
        - Financial Impact: Estimated at {rand_money(5_000, 500_000)} including investigation and remediation costs

        TIMELINE
        - {rand_date(days_back=90)} 09:{random.randint(10, 59)} -- Initial detection via automated alert
        - {rand_date(days_back=90)} 10:{random.randint(10, 59)} -- SOC team engaged, triage initiated
        - {rand_date(days_back=89)} -- Containment actions completed
        - {rand_date(days_back=85)} -- Root cause analysis completed
        - {rand_date(days_back=80)} -- Remediation actions implemented

        ROOT CAUSE
        {random.choice(["Employee clicked malicious link in sophisticated phishing email that bypassed email security controls.", "Compromised third-party credentials used to access internal systems via VPN.", "Unpatched vulnerability in web application exploited by external threat actor.", "Misconfigured cloud storage bucket exposed to public internet."])}

        CORRECTIVE ACTIONS
        1. Implemented additional email filtering rules targeting identified attack patterns.
        2. Forced password reset for all potentially affected accounts.
        3. Deployed enhanced endpoint detection rules for identified threat indicators.
        4. Conducted targeted security awareness training for affected business unit.

        Report prepared by: {rand_name()}, CISO
    """))
ENT_TEMPLATES.append(_ent_security_incident)

def _ent_budget_forecast():
    fiscal = rand_fiscal_year()
    total = random.randint(20, 500)
    return ("budget_forecast", f"Annual Budget Forecast -- {fiscal}", textwrap.dedent(f"""\
        ANNUAL BUDGET FORECAST
        Classification: {random.choice(CLASSIFICATIONS_ENT)}
        Fiscal Year: {fiscal}
        Date: {rand_date()}
        Prepared by: {rand_name()}, Director of Finance
        Organization: {random.choice(ORGS_ENT)}

        BUDGET SUMMARY
        Total Budget: ${total}M
        Revenue Target: ${random.randint(total, total * 3)}M
        Capital Expenditure: ${random.randint(1, total // 3)}M
        Operating Expense: ${random.randint(total // 2, total)}M
        Contingency Reserve: ${random.randint(1, total // 10)}M

        DEPARTMENT ALLOCATIONS
        Engineering: ${random.randint(5, total // 3)}M ({rand_pct(20, 40)} of total)
        Sales & Marketing: ${random.randint(3, total // 4)}M ({rand_pct(15, 30)} of total)
        Operations: ${random.randint(2, total // 4)}M ({rand_pct(10, 25)} of total)
        General & Administrative: ${random.randint(1, total // 5)}M ({rand_pct(5, 15)} of total)
        Research & Development: ${random.randint(1, total // 4)}M ({rand_pct(8, 20)} of total)

        KEY ASSUMPTIONS
        - Revenue growth rate: {rand_pct(5, 25)} year-over-year
        - Headcount growth: {random.randint(10, 200)} net new hires
        - Average salary increase: {rand_pct(3, 8)} per employee
        - Infrastructure cost reduction from cloud migration: {rand_pct(10, 30)}
        - Foreign exchange impact: {random.choice(["+", "-"])}{rand_pct(1, 5)} on international revenue

        RISK FACTORS
        1. Economic uncertainty may impact client spending decisions.
        2. Competitive hiring market could drive compensation costs above forecast.
        3. Regulatory changes may require unbudgeted compliance investments.

        APPROVAL CHAIN
        {rand_name()}, CFO -- Approved
        {rand_name()}, CEO -- Approved
        Board of Directors -- {random.choice(["Approved", "Pending Review"])}
    """))
ENT_TEMPLATES.append(_ent_budget_forecast)


# ─────────────────────────────────────────────────────────────────────────────
# GOLDEN ENTERPRISE DOCUMENTS -- deterministic content for test queries
# ─────────────────────────────────────────────────────────────────────────────
# These documents contain fixed facts that answer the specific test queries
# in the Enterprise Test Plan.  They are always generated as .txt and are
# NOT randomised, so the test harness can verify deterministic answers.
#
# Golden doc source files live in golden_docs/ alongside this script.
# ─────────────────────────────────────────────────────────────────────────────

GOLDEN_DOCS_DIR = Path(__file__).resolve().parent / "golden_docs"


# ─────────────────────────────────────────────────────────────────────────────
# Format writers
# ─────────────────────────────────────────────────────────────────────────────
def write_txt(path: Path, content: str) -> None:
    path.write_text(content, encoding="utf-8")

def write_md(path: Path, title: str, content: str) -> None:
    md = f"# {title}\n\n"
    for line in content.splitlines():
        stripped = line.strip()
        if stripped and stripped == stripped.upper() and len(stripped) > 5 and not stripped.startswith(("-", "1", "2", "3", "4", "5")):
            md += f"\n## {stripped.title()}\n\n"
        elif stripped:
            md += stripped + "\n\n"
    path.write_text(md, encoding="utf-8")

def write_html(path: Path, title: str, content: str) -> None:
    body = ""
    for line in content.splitlines():
        stripped = line.strip()
        if stripped and stripped == stripped.upper() and len(stripped) > 5 and not stripped.startswith(("-", "1", "2", "3", "4", "5")):
            body += f"<h2>{stripped.title()}</h2>\n"
        elif stripped:
            body += f"<p>{stripped}</p>\n"
    html = f"<!DOCTYPE html>\n<html><head><title>{title}</title></head>\n<body>\n<h1>{title}</h1>\n{body}</body></html>"
    path.write_text(html, encoding="utf-8")

def write_csv(path: Path, content: str) -> None:
    lines = [l.strip() for l in content.splitlines() if l.strip() and ":" in l]
    with path.open("w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["Field", "Value"])
        for line in lines[:50]:
            parts = line.split(":", 1)
            if len(parts) == 2:
                writer.writerow([parts[0].strip(), parts[1].strip()])

def write_json_doc(path: Path, title: str, content: str, sector: str) -> None:
    payload = {
        "doc_id": rand_id("DOC"),
        "sector": sector,
        "title": title,
        "created_at": datetime.now(timezone.utc).isoformat().replace("+00:00", "Z"),
        "content": content.splitlines()[:120],
    }
    path.write_text(json.dumps(payload, indent=2), encoding="utf-8")

def write_ndjson(path: Path, content: str, sector: str) -> None:
    lines = [l.strip() for l in content.splitlines() if l.strip()][:120]
    with path.open("w", encoding="utf-8") as f:
        f.write(json.dumps({"doc_id": rand_id("DOC"), "sector": sector}) + "\n")
        for line in lines:
            f.write(json.dumps({"line": line}) + "\n")

def write_log(path: Path, content: str) -> None:
    lines = [l.strip() for l in content.splitlines() if l.strip()][:120]
    now = datetime.now(timezone.utc)
    with path.open("w", encoding="utf-8") as f:
        for i, line in enumerate(lines):
            ts = (now + timedelta(seconds=i)).replace(microsecond=0).isoformat() + "Z"
            level = random.choice(["INFO", "INFO", "INFO", "WARN", "DEBUG"])
            f.write(f"{ts} {level} {line}\n")

def write_docx(path: Path, title: str, content: str) -> None:
    if not HAS_DOCX:
        write_txt(path.with_suffix(".txt"), content)
        return
    doc = DocxDocument()
    doc.add_heading(title, level=1)
    for line in content.splitlines()[:100]:
        stripped = line.strip()
        if stripped:
            if stripped == stripped.upper() and len(stripped) > 5 and not stripped.startswith(("-", "1", "2", "3", "4", "5")):
                doc.add_heading(stripped.title(), level=2)
            else:
                doc.add_paragraph(stripped)
    doc.save(path)

def write_pptx(path: Path, title: str, content: str) -> None:
    if not HAS_PPTX:
        write_txt(path.with_suffix(".txt"), content)
        return
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    bullets = [l.strip() for l in content.splitlines() if l.strip() and not l.strip() == l.strip().upper()][:6]
    slide.placeholders[1].text = "\n".join(bullets) if bullets else "See attached documentation."
    prs.save(path)

def write_xlsx(path: Path, title: str, content: str) -> None:
    if not HAS_XLSX:
        write_csv(path.with_suffix(".csv"), content)
        return
    wb = Workbook()
    ws = wb.active
    ws.title = "Report"
    ws.append([title])
    ws.append([])
    for line in content.splitlines()[:200]:
        stripped = line.strip()
        if stripped:
            if ":" in stripped:
                parts = stripped.split(":", 1)
                ws.append([parts[0].strip(), parts[1].strip()])
            else:
                ws.append([stripped])
    wb.save(path)

def _sanitize_latin1(text: str) -> str:
    """Replace non-latin-1 characters with ASCII equivalents for fpdf2."""
    return (text
            .replace("\u2014", "--")   # em-dash
            .replace("\u2013", "-")    # en-dash
            .replace("\u2018", "'")    # left single quote
            .replace("\u2019", "'")    # right single quote
            .replace("\u201c", '"')    # left double quote
            .replace("\u201d", '"')    # right double quote
            .replace("\u2026", "...")  # ellipsis
            .replace("\u2022", "*")   # bullet
            .encode("latin-1", errors="replace").decode("latin-1"))

def write_pdf(path: Path, title: str, content: str) -> None:
    if not HAS_PDF:
        write_txt(path.with_suffix(".txt"), content)
        return
    from fpdf.enums import XPos, YPos
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 14)
    pdf.cell(0, 10, _sanitize_latin1(title[:80]), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.set_font("Helvetica", "", 10)
    for line in content.splitlines()[:150]:
        stripped = line.strip()
        if stripped:
            pdf.multi_cell(0, 5, _sanitize_latin1(stripped))
            pdf.ln(1)
    pdf.output(str(path))


# ─────────────────────────────────────────────────────────────────────────────
# Format selection (weighted)
# ─────────────────────────────────────────────────────────────────────────────
FORMAT_WEIGHTS: list[tuple[str, float]] = [
    (".txt",    0.40),
    (".docx",   0.12),
    (".xlsx",   0.06),
    (".csv",    0.04),
    (".html",   0.08),
    (".md",     0.05),
    (".pptx",   0.03),
    (".json",   0.04),
    (".ndjson",  0.02),
    (".log",    0.01),
]
if HAS_PDF:
    FORMAT_WEIGHTS.insert(1, (".pdf", 0.15))
    # Rebalance: reduce .txt to 0.25
    FORMAT_WEIGHTS[0] = (".txt", 0.25)

FORMATS = [f for f, _ in FORMAT_WEIGHTS]
WEIGHTS = [w for _, w in FORMAT_WEIGHTS]


def pick_format() -> str:
    return random.choices(FORMATS, weights=WEIGHTS, k=1)[0]


def write_document(path: Path, title: str, content: str, sector: str, fmt: str) -> None:
    """Write content in the specified format."""
    if fmt == ".txt":
        write_txt(path, content)
    elif fmt == ".md":
        write_md(path, title, content)
    elif fmt == ".html":
        write_html(path, title, content)
    elif fmt == ".csv":
        write_csv(path, content)
    elif fmt == ".json":
        write_json_doc(path, title, content, sector)
    elif fmt == ".ndjson":
        write_ndjson(path, content, sector)
    elif fmt == ".log":
        write_log(path, content)
    elif fmt == ".docx":
        write_docx(path, title, content)
    elif fmt == ".pptx":
        write_pptx(path, title, content)
    elif fmt == ".xlsx":
        write_xlsx(path, title, content)
    elif fmt == ".pdf":
        write_pdf(path, title, content)
    else:
        write_txt(path, content)


# ─────────────────────────────────────────────────────────────────────────────
# Golden document writer
# ─────────────────────────────────────────────────────────────────────────────
def generate_golden_docs(base: Path) -> int:
    """Copy deterministic golden enterprise documents to Enterprise/ subdirectory.

    Golden doc source files live in ``golden_docs/`` alongside this script.
    They are always .txt files with fixed content and are copied verbatim
    BEFORE the randomised corpus so they exist regardless of --count.
    Returns number of golden docs written.
    """
    sector_dir = base / "Enterprise"
    sector_dir.mkdir(parents=True, exist_ok=True)

    if not GOLDEN_DOCS_DIR.is_dir():
        print(f"  WARNING: golden_docs/ directory not found at {GOLDEN_DOCS_DIR}")
        return 0

    count = 0
    for src_file in sorted(GOLDEN_DOCS_DIR.glob("*.txt")):
        dst = sector_dir / src_file.name
        shutil.copy2(src_file, dst)
        count += 1
        print(f"    [golden] {dst.name}")
    return count


# ─────────────────────────────────────────────────────────────────────────────
# Main generation loop
# ─────────────────────────────────────────────────────────────────────────────
SECTOR_TEMPLATES = {
    "GOVERNMENT": GOV_TEMPLATES,
    "MEDICAL":    MED_TEMPLATES,
    "ENTERPRISE": ENT_TEMPLATES,
}

SECTOR_DIRS = {
    "GOVERNMENT": "Government",
    "MEDICAL":    "Medical",
    "ENTERPRISE": "Enterprise",
}


def main():
    parser = argparse.ArgumentParser(description="Generate enterprise corpus for Sentinel RAG testing")
    parser.add_argument("--output", default=r"D:\Corpus", help="Output directory (default: D:\\Corpus)")
    parser.add_argument("--count", type=int, default=3000, help="Total number of documents (split evenly across sectors)")
    parser.add_argument("--seed", type=int, default=42, help="Random seed for reproducibility")
    parser.add_argument("--golden-only", action="store_true", help="Only generate golden test documents")
    args = parser.parse_args()

    random.seed(args.seed)
    base = Path(args.output)

    # Always generate golden docs first
    print("Generating golden enterprise documents (deterministic test fixtures)...")
    golden_count = generate_golden_docs(base)
    print(f"  {golden_count} golden documents written.\n")

    if args.golden_only:
        print("--golden-only specified, skipping randomized corpus generation.")
        return

    per_sector = args.count // len(SECTOR_TEMPLATES)

    print(f"Generating {args.count} randomized documents ({per_sector} per sector) to {base}")
    print(f"PDF support: {'YES' if HAS_PDF else 'NO (install fpdf2 for PDF output)'}")
    print(f"DOCX support: {'YES' if HAS_DOCX else 'NO'}")
    print(f"PPTX support: {'YES' if HAS_PPTX else 'NO'}")
    print(f"XLSX support: {'YES' if HAS_XLSX else 'NO'}")
    print()

    total = 0
    format_counts: dict[str, int] = {}

    for sector, templates in SECTOR_TEMPLATES.items():
        sector_dir = base / SECTOR_DIRS[sector]
        sector_dir.mkdir(parents=True, exist_ok=True)

        print(f"  {sector}: generating {per_sector} documents...")

        for i in range(1, per_sector + 1):
            template_fn = random.choice(templates)
            category, title, content = template_fn()

            # Inject PII into ~5% of docs
            content = inject_pii(content, sector)

            # Pick output format
            fmt = pick_format()
            format_counts[fmt] = format_counts.get(fmt, 0) + 1

            # Build filename
            prefix = {"GOVERNMENT": "gov", "MEDICAL": "med", "ENTERPRISE": "ent"}[sector]
            filename = f"{prefix}_{category}_{i:04d}{fmt}"
            filepath = sector_dir / filename

            write_document(filepath, title, content, sector, fmt)
            total += 1

            if i % 250 == 0:
                print(f"    ... {i}/{per_sector}")

    print()
    total_with_golden = total + golden_count
    print(f"Done. Generated {total} randomized + {golden_count} golden = {total_with_golden} total files.")
    print(f"Format distribution (randomized only):")
    for fmt, count in sorted(format_counts.items(), key=lambda x: -x[1]):
        print(f"  {fmt:>8}: {count:>5} ({count * 100 / total:.0f}%)")


if __name__ == "__main__":
    main()
